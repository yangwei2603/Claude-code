#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
内容分析模块 — 文件内容提取与文件名分析
从 organizer.py 拆分出来：ContentExtractor、FilenameAnalyzer
"""

import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Optional

# 魔法数字
MAX_EXTRACTED_TEXT_CHARS = 50000
HASH_CHUNK_SIZE = 65536


class ContentExtractor:
    """提取文件内容的工具类"""

    def __init__(self, max_size_mb: int = 10):
        self.max_size_bytes = max_size_mb * 1024 * 1024

    def can_extract(self, filepath: Path) -> bool:
        """检查是否支持提取该文件的内容"""
        ext = filepath.suffix.lower()
        supported = {'.txt', '.md', '.docx', '.pdf', '.doc', '.rtf', '.csv',
                     '.xlsx', '.xls', '.pptx', '.ppt', '.json', '.xml',
                     '.html', '.htm', '.py', '.js', '.java', '.sql', '.yaml', '.yml'}

        if ext not in supported:
            return False

        try:
            if filepath.stat().st_size > self.max_size_bytes:
                return False
        except Exception:
            return False

        return True

    def extract(self, filepath: Path) -> str:
        """提取文件内容，返回文本字符串"""
        ext = filepath.suffix.lower()

        try:
            if ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm', '.py', '.js', '.java', '.sql', '.yaml', '.yml']:
                return self._extract_text(filepath)
            elif ext == '.docx':
                return self._extract_docx(filepath)
            elif ext == '.pdf':
                return self._extract_pdf(filepath)
            elif ext in ['.xlsx', '.xls']:
                return self._extract_excel(filepath)
            elif ext in ['.pptx', '.ppt']:
                return self._extract_ppt(filepath)
            elif ext == '.doc':
                return self._extract_doc(filepath)
            elif ext == '.rtf':
                return self._extract_rtf(filepath)
        except Exception:
            return ""

        return ""

    def _extract_text(self, filepath: Path) -> str:
        """提取纯文本文件内容"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding, errors='ignore') as f:
                    return f.read()[:MAX_EXTRACTED_TEXT_CHARS]
            except Exception:
                continue
        return ""

    def _extract_docx(self, filepath: Path) -> str:
        """提取 Word 文档内容"""
        try:
            from docx import Document
            doc = Document(filepath)
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text.append(cell.text)
            return '\n'.join(text)[:MAX_EXTRACTED_TEXT_CHARS]
        except ImportError:
            return self._extract_docx_xml(filepath)
        except Exception:
            return ""

    def _extract_docx_xml(self, filepath: Path) -> str:
        """通过 XML 提取 docx 内容（备用方案）"""
        try:
            with zipfile.ZipFile(filepath) as z:
                xml_content = z.read('word/document.xml')
                tree = ET.fromstring(xml_content)
                texts = [elem.text for elem in tree.iter() if elem.text]
                return ' '.join(texts)[:MAX_EXTRACTED_TEXT_CHARS]
        except Exception:
            return ""

    def _extract_pdf(self, filepath: Path) -> str:
        """提取 PDF 内容"""
        try:
            import PyPDF2
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = [page.extract_text() or "" for page in reader.pages]
                return '\n'.join(text)[:MAX_EXTRACTED_TEXT_CHARS]
        except ImportError:
            return ""
        except Exception:
            return ""

    def _extract_excel(self, filepath: Path) -> str:
        """提取 Excel 内容"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text.append(str(cell.value))
            return ' '.join(text)[:MAX_EXTRACTED_TEXT_CHARS]
        except ImportError:
            return ""
        except Exception:
            return ""

    def _extract_ppt(self, filepath: Path) -> str:
        """提取 PPT 内容"""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = [shape.text for slide in prs.slides if hasattr(shape, "text") for shape in slide.shapes]
            return '\n'.join(text)[:MAX_EXTRACTED_TEXT_CHARS]
        except ImportError:
            return ""
        except Exception:
            return ""

    def _extract_doc(self, filepath: Path) -> str:
        """提取旧版 Word 文档内容（.doc）"""
        # .doc 格式较复杂，需要额外库支持
        return ""

    def _extract_rtf(self, filepath: Path) -> str:
        """提取 RTF 内容"""
        try:
            from striprtf.striprtf import rtf_to_text
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return rtf_to_text(f.read())[:MAX_EXTRACTED_TEXT_CHARS]
        except ImportError:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            cleaned = re.sub(r'\\[a-z]+\d*\s?', ' ', content)
            cleaned = re.sub(r'[{}]', '', cleaned)
            return cleaned[:MAX_EXTRACTED_TEXT_CHARS]
        except Exception:
            return ""


class FilenameAnalyzer:
    """分析文件名提取结构化信息"""

    # 日期模式
    DATE_PATTERNS = [
        r'\d{4}[-年]\d{1,2}[-月]\d{1,2}[日]?',  # 2024-01-15, 2024年01月15日
        r'\d{4}\d{2}\d{2}',  # 20240115
        r'\d{2}[-/]\d{1,2}[-/]\d{1,2}',  # 24-01-15
    ]

    # 项目编号模式
    PROJECT_PATTERNS = [
        r'[A-Z]{2,4}-\d{3,}',  # ERP-001, DIG-2024-001
        r'PRJ\d{3,}',  # PRJ001
        r'项目\d+',  # 项目001
    ]

    # 版本号模式
    VERSION_PATTERNS = [
        r'V\d+\.\d+',  # V1.0, V2.5
        r'v\d+\.\d+',  # v1.0
        r'\d+\.\d+\.\d+',  # 1.0.0
    ]

    # 文档类型关键词
    DOC_TYPE_KEYWORDS = {
        '需求': ['需求', 'BRD', 'PRD', 'requirement'],
        '设计': ['设计', 'design', '架构', 'architecture'],
        '测试': ['测试', 'test', '用例', 'case'],
        '报告': ['报告', 'report', '分析', 'analysis'],
        '手册': ['手册', 'manual', '指南', 'guide'],
        '计划': ['计划', 'plan', 'schedule', '排期'],
        '总结': ['总结', 'summary', '复盘', 'review'],
        '会议纪要': ['纪要', 'minutes', 'meeting'],
    }

    def analyze(self, filename: str) -> dict:
        """分析文件名，返回结构化信息"""
        result = {
            'dates': [],
            'project_codes': [],
            'versions': [],
            'doc_types': [],
            'keywords': []
        }

        for pattern in self.DATE_PATTERNS:
            result['dates'].extend(re.findall(pattern, filename))

        for pattern in self.PROJECT_PATTERNS:
            result['project_codes'].extend(re.findall(pattern, filename, re.IGNORECASE))

        for pattern in self.VERSION_PATTERNS:
            result['versions'].extend(re.findall(pattern, filename))

        for doc_type, keywords in self.DOC_TYPE_KEYWORDS.items():
            for kw in keywords:
                if kw.lower() in filename.lower():
                    result['doc_types'].append(doc_type)
                    break

        return result
