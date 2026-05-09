"""
内容提取器 - 从文件中提取文本内容
"""

import re
import zipfile
import logging
from pathlib import Path
from typing import Optional

from .constants import DEFAULT_MAX_CONTENT_SIZE_MB, BYTES_PER_MB

logger = logging.getLogger("content_extractor")


class ContentExtractor:
    """提取文件内容的工具类"""

    def __init__(self, max_size_mb: int = DEFAULT_MAX_CONTENT_SIZE_MB):
        self.max_size_bytes = max_size_mb * BYTES_PER_MB
        self.supported_extensions = {
            '.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm',
            '.py', '.js', '.java', '.sql', '.yaml', '.yml',
            '.docx', '.pdf', '.xlsx', '.xls', '.pptx', '.ppt', '.rtf'
        }

    def can_extract(self, filepath: Path) -> bool:
        """检查是否支持提取该文件的内容"""
        ext = filepath.suffix.lower()

        if ext not in self.supported_extensions:
            return False

        try:
            if filepath.stat().st_size > self.max_size_bytes:
                return False
        except OSError:
            return False

        return True

    def extract(self, filepath: Path) -> str:
        """提取文件内容，返回文本字符串"""
        ext = filepath.suffix.lower()

        if not self.can_extract(filepath):
            return ""

        try:
            if ext in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm',
                       '.py', '.js', '.java', '.sql', '.yaml', '.yml']:
                return self._extract_text(filepath)
            elif ext == '.docx':
                return self._extract_docx(filepath)
            elif ext == '.pdf':
                return self._extract_pdf(filepath)
            elif ext in ['.xlsx', '.xls']:
                return self._extract_excel(filepath)
            elif ext in ['.pptx', '.ppt']:
                return self._extract_ppt(filepath)
            elif ext == '.rtf':
                return self._extract_rtf(filepath)
        except Exception as e:
            logger.debug(f"内容提取失败（{filepath.name}）: {e}")
            return ""

        return ""

    def _extract_text(self, filepath: Path) -> str:
        """提取纯文本文件内容"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding, errors='ignore') as f:
                    content = f.read()
                    return content[:50000]  # 最多返回5万字符
            except Exception:
                continue
        return ""

    def _extract_docx(self, filepath: Path) -> str:
        """提取 Word 文档内容"""
        try:
            from docx import Document
            doc = Document(filepath)
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text.append(cell.text)
            return '\n'.join(text)[:50000]
        except ImportError:
            return self._extract_docx_xml(filepath)
        except Exception:
            return ""

    def _extract_docx_xml(self, filepath: Path) -> str:
        """通过 XML 提取 docx 内容（备用方案）"""
        try:
            with zipfile.ZipFile(filepath) as z:
                xml_content = z.read('word/document.xml')
                import xml.etree.ElementTree as ET
                tree = ET.fromstring(xml_content)
                texts = []
                for elem in tree.iter():
                    if elem.text:
                        texts.append(elem.text)
                return ' '.join(texts)[:50000]
        except Exception:
            return ""

    def _extract_pdf(self, filepath: Path) -> str:
        """提取 PDF 内容"""
        try:
            import PyPDF2
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = []
                for page in reader.pages:
                    text.append(page.extract_text() or "")
                return '\n'.join(text)[:50000]
        except ImportError:
            return ""
        except Exception:
            return ""

    def _extract_excel(self, filepath: Path) -> str:
        """提取 Excel 内容"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            text.append(str(cell.value))
            return ' '.join(text)[:50000]
        except ImportError:
            return ""
        except Exception:
            return ""

    def _extract_ppt(self, filepath: Path) -> str:
        """提取 PPT 内容"""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)[:50000]
        except ImportError:
            return ""
        except Exception:
            return ""

    def _extract_rtf(self, filepath: Path) -> str:
        """提取 RTF 内容"""
        try:
            from striprtf.striprtf import rtf_to_text
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            return rtf_to_text(rtf_content)[:50000]
        except ImportError:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            cleaned = re.sub(r'\\[a-z]+\d*\s?', ' ', content)
            cleaned = re.sub(r'[{}]', '', cleaned)
            return cleaned[:50000]
        except Exception:
            return ""
